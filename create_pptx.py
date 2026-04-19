from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Add a blank slide
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)

# Background color (Dark Navy)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 23, 42)

def add_textbox(slide, text, left, top, width, height, font_size=18, color=(248, 250, 252), bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*color)
    p.font.bold = bold
    p.alignment = align
    return txBox

# Title
add_textbox(slide, "Meta-Analysis: Methods & Pipeline Framework", 0.5, 0.4, 9, 1, font_size=36, color=(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)

# Subtitle / Goal
add_textbox(slide, "GOAL: Integrated microbiome analysis across datasets and sequencing approaches", 0.5, 1.3, 9, 0.5, font_size=14, color=(129, 140, 248), bold=True, align=PP_ALIGN.CENTER)

# Framework description
add_textbox(slide, "Used a multi-pipeline framework to improve comparability and reproducibility.", 1, 1.8, 8, 0.4, font_size=16, color=(148, 163, 184), align=PP_ALIGN.CENTER)

# Pipelines Grid (Simulation)
# Pipeline 1
add_textbox(slide, "PIPELINE 1: Greengenes2 + Woltka", 0.5, 2.5, 2.8, 0.5, font_size=14, color=(248, 250, 252), bold=True)
add_textbox(slide, "• Shared GG2 taxonomy\n• 16S: QIIME2/DADA2\n• Shotgun: WoLr2/Woltka", 0.5, 2.9, 2.8, 1.5, font_size=12, color=(148, 163, 184))

# Pipeline 2
add_textbox(slide, "PIPELINE 2: RefSeq + Kraken/Bracken", 3.6, 2.5, 2.8, 0.5, font_size=14, color=(248, 250, 252), bold=True)
add_textbox(slide, "• NCBI RefSeq DB\n• 16S: RefSeq Classifier\n• Shotgun: Kraken2/Bracken", 3.6, 2.9, 2.8, 1.5, font_size=12, color=(148, 163, 184))

# Pipeline 3
add_textbox(slide, "PIPELINE 3: 16S Extraction (SortMeRNA)", 6.7, 2.5, 2.8, 0.5, font_size=14, color=(248, 250, 252), bold=True)
add_textbox(slide, "• Marker Gene Isolation\n• Shotgun 16S Extraction\n• Harmonized Processing", 6.7, 2.9, 2.8, 1.5, font_size=12, color=(148, 163, 184))

# GitHub Link
add_textbox(slide, "GitHub Repository: https://github.com/gio9024/16S-and-SM-data-integration", 0.5, 6.5, 9, 0.5, font_size=12, color=(255, 255, 255), align=PP_ALIGN.CENTER)

# Save
prs.save('project_overview.pptx')
print("Successfully created project_overview.pptx")
