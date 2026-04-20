"""
Generate summary_slide.pptx — metagenomics pipeline project summary
White/light theme. One slide: 3 pipelines | 4 databases | 2 completed datasets
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Widescreen slide (13.33" × 7.5") ───────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ─── Color palette (LIGHT THEME) ─────────────────────────────────────────────
C_BG        = RGBColor(0xFF, 0xFF, 0xFF)   # white background
C_SURFACE   = RGBColor(0xF8, 0xFA, 0xFC)   # card surface
C_SURFACE2  = RGBColor(0xF1, 0xF5, 0xF9)   # stat box surface
C_BORDER    = RGBColor(0xE2, 0xE8, 0xF0)   # card border
C_TEXT      = RGBColor(0x0F, 0x17, 0x2A)   # main text (dark navy)
C_MUTED     = RGBColor(0x64, 0x74, 0x8B)   # secondary text

# Pipeline accent colors (kept vivid for contrast on white)
C_P1        = RGBColor(0x63, 0x66, 0xF1)   # indigo
C_P1L       = RGBColor(0x43, 0x38, 0xCA)   # darker indigo for text on light bg
C_P1BG      = RGBColor(0xED, 0xE9, 0xFE)   # very light indigo
C_P1LINE    = RGBColor(0xA5, 0xB4, 0xFC)

C_P2        = RGBColor(0x0E, 0xA5, 0xE9)   # sky blue
C_P2L       = RGBColor(0x03, 0x69, 0xA1)   # darker sky for text
C_P2BG      = RGBColor(0xE0, 0xF2, 0xFE)
C_P2LINE    = RGBColor(0x7D, 0xD3, 0xFC)

C_P3        = RGBColor(0x10, 0xB9, 0x81)   # emerald
C_P3L       = RGBColor(0x06, 0x5F, 0x46)   # darker emerald for text
C_P3BG      = RGBColor(0xD1, 0xFA, 0xE5)
C_P3LINE    = RGBColor(0x6E, 0xE7, 0xB7)

C_GOLD      = RGBColor(0xD9, 0x77, 0x06)   # amber (darker for light bg)
C_GOLDBG    = RGBColor(0xFF, 0xF7, 0xED)
C_GOLDLINE  = RGBColor(0xFB, 0xBF, 0x24)

C_PINK      = RGBColor(0xBE, 0x18, 0x5D)   # pink (darker)
C_PINKBG    = RGBColor(0xFD, 0xF2, 0xF8)
C_PINKLINE  = RGBColor(0xF9, 0xA8, 0xD4)

C_GREEN_OK  = RGBColor(0x05, 0x96, 0x69)
C_GREEN_BG  = RGBColor(0xEC, 0xFD, 0xF5)

# ─── Helpers ─────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0.75)):
    shape = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=Pt(11), bold=False, color=None,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_TEXT
    return txb


def add_rounded_rect(slide, l, t, w, h, fill=None, line=None,
                     line_w=Pt(0.75), radius=0.08):
    shape = slide.shapes.add_shape(
        5, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.adjustments[0] = radius
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def section_label(text, l, t):
    add_text_box(slide, text.upper(), l, t, 8, 0.18,
                 font_size=Pt(6.5), bold=True, color=C_MUTED)


# ─── SLIDE BACKGROUND ────────────────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = C_BG

# ─── HEADER ──────────────────────────────────────────────────────────────────
# Title
title_txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.72))
title_txb.word_wrap = False
tf = title_txb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "16S & Shotgun Metagenomics Integration — Project Summary"
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = C_TEXT

# Subtitle
add_text_box(slide,
    "Three parallel pipelines benchmarked across two QC datasets, comparing 16S amplicon "
    "and whole-genome shotgun classification with harmonized taxonomy.",
    1.0, 0.92, 11.33, 0.32,
    font_size=Pt(9), color=C_MUTED, align=PP_ALIGN.CENTER)

# Thin divider line
add_rect(slide, 0.5, 1.25, 12.33, 0.012, fill=C_BORDER)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — THREE PIPELINES
# ═══════════════════════════════════════════════════════════════════════════════
section_label("Three Analysis Pipelines", 0.5, 1.3)

PIPE_CONFIGS = [
    {
        "num":       "Pipeline 1",
        "title":     "Greengenes2 + Woltka",
        "accent":    C_P1,
        "accentL":   C_P1L,
        "accent_bg": C_P1BG,
        "accentLn":  C_P1LINE,
        "16s_flow":  "DADA2  →  GG2 V4 NB Classifier  →  Genus Table",
        "mgs_flow":  "Bowtie2 / WoLr2  →  Woltka OGU  →  Genus Table",
        "advantage": "Both sides share the same Greengenes2 (2024.09) taxonomy tree — "
                     "enabling direct label-level comparison. Uses GTDB taxonomy.",
        "col": 0
    },
    {
        "num":       "Pipeline 2",
        "title":     "RefSeq + Kraken2/Bracken",
        "accent":    C_P2,
        "accentL":   C_P2L,
        "accent_bg": C_P2BG,
        "accentLn":  C_P2LINE,
        "16s_flow":  "DADA2  →  RefSeq 16S V4 Classifier  →  Genus Table",
        "mgs_flow":  "Kraken2 (full RefSeq)  →  Bracken  →  Genus Table",
        "advantage": "Full NCBI RefSeq on both sides. Kraken2+Bracken: fast k-mer "
                     "classification. Best specificity: fewest FP among all MGS approaches.",
        "col": 1
    },
    {
        "num":       "Pipeline 3",
        "title":     "SortMeRNA → Kraken2/Bracken",
        "accent":    C_P3,
        "accentL":   C_P3L,
        "accent_bg": C_P3BG,
        "accentLn":  C_P3LINE,
        "16s_flow":  "DADA2  →  RefSeq 16S V4 Classifier  →  Genus Table",
        "mgs_flow":  "SortMeRNA 16S extract  →  Kraken2 + Bracken  →  Genus Table",
        "advantage": "Extracts only 16S rRNA reads from WGS before classification — "
                     "isolating the same marker gene as amplicon data for maximum parity.",
        "col": 2
    }
]

PIPE_L = 0.5
PIPE_T = 1.48
PIPE_W = 4.05
PIPE_H = 1.9
PIPE_GAP = 0.115

for cfg in PIPE_CONFIGS:
    l = PIPE_L + cfg["col"] * (PIPE_W + PIPE_GAP)
    t = PIPE_T

    # Card background
    add_rounded_rect(slide, l, t, PIPE_W, PIPE_H,
                     fill=C_SURFACE, line=C_BORDER, line_w=Pt(0.75))

    # Accent top bar
    add_rounded_rect(slide, l, t, PIPE_W, 0.07,
                     fill=cfg["accent"], radius=0.12)

    # Badge
    add_rounded_rect(slide, l + 0.15, t + 0.12, 0.9, 0.21,
                     fill=cfg["accent_bg"],
                     line=cfg["accentLn"], line_w=Pt(0.5), radius=0.2)
    add_text_box(slide, cfg["num"].upper(), l + 0.15, t + 0.13, 0.9, 0.19,
                 font_size=Pt(6), bold=True, color=cfg["accentL"],
                 align=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, cfg["title"], l + 0.15, t + 0.36, PIPE_W - 0.3, 0.26,
                 font_size=Pt(10.5), bold=True, color=C_TEXT)

    # 16S flow
    add_text_box(slide, "16S", l + 0.15, t + 0.64, 0.3, 0.18,
                 font_size=Pt(6.5), bold=True, color=cfg["accentL"])
    add_text_box(slide, cfg["16s_flow"], l + 0.48, t + 0.64, PIPE_W - 0.63, 0.2,
                 font_size=Pt(7.5), color=C_TEXT)

    # MGS flow
    add_text_box(slide, "MGS", l + 0.15, t + 0.87, 0.33, 0.18,
                 font_size=Pt(6.5), bold=True, color=cfg["accentL"])
    add_text_box(slide, cfg["mgs_flow"], l + 0.48, t + 0.87, PIPE_W - 0.63, 0.22,
                 font_size=Pt(7.5), color=C_TEXT)

    # Divider
    add_rect(slide, l + 0.15, t + 1.14, PIPE_W - 0.3, 0.008, fill=C_BORDER)

    # Advantage text
    add_text_box(slide, cfg["advantage"], l + 0.15, t + 1.18, PIPE_W - 0.3, 0.66,
                 font_size=Pt(7.5), color=C_MUTED, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 2 — DATABASES
# ═══════════════════════════════════════════════════════════════════════════════
DB_T = 3.48
section_label("Reference Databases", 0.5, DB_T - 0.17)

DB_CARD_L = 0.5
DB_CARD_T = DB_T
DB_CARD_W = 12.33
DB_CARD_H = 1.18

add_rounded_rect(slide, DB_CARD_L, DB_CARD_T, DB_CARD_W, DB_CARD_H,
                 fill=C_SURFACE, line=C_BORDER, line_w=Pt(0.75))

DATABASES = [
    {
        "icon": "🌿",
        "name": "Greengenes2",
        "ver":  "v2024.09 · GTDB taxonomy",
        "desc": "Phylogenomic reference integrating full-length 16S with WGS. "
                "Single tree for both QIIME2 and Woltka.",
        "tags": [("P1 16S", C_P1L, C_P1BG, C_P1LINE),
                 ("P1 MGS (tax)", C_P1L, C_P1BG, C_P1LINE)],
    },
    {
        "icon": "🌐",
        "name": "WoLr2 (Web of Life r2)",
        "ver":  "Bowtie2 index · Prokaryotic genomes",
        "desc": "Prokaryotic genome DB for WGS read alignment. Reads mapped "
                "via Bowtie2, then classified by Woltka using GG2 taxonomy.",
        "tags": [("P1 MGS (align)", C_P1L, C_P1BG, C_P1LINE)],
    },
    {
        "icon": "🧬",
        "name": "NCBI RefSeq (full genome)",
        "ver":  "Build 2024-01-12 · NCBI taxonomy",
        "desc": "Comprehensive curated genomes: bacteria, archaea, fungi, viruses, "
                "human. Kraken2 standard hash 72 GB.",
        "tags": [("P2 MGS", C_P2L, C_P2BG, C_P2LINE),
                 ("P3 MGS", C_P3L, C_P3BG, C_P3LINE)],
    },
    {
        "icon": "🔬",
        "name": "NCBI 16S RefSeq",
        "ver":  "Nov 2024 · 26,244 dereplicated seqs",
        "desc": "Curated 16S rRNA sequences. V4-region NB classifier for amplicons; "
                "full-length NB for extracted reads.",
        "tags": [("P2 16S", C_P2L, C_P2BG, C_P2LINE),
                 ("P3 16S", C_P3L, C_P3BG, C_P3LINE),
                 ("P3 MGS", C_P3L, C_P3BG, C_P3LINE)],
    },
    {
        "icon": "🎯",
        "name": "SortMeRNA rRNA (SILVA)",
        "ver":  "v4.3 · SILVA-based rRNA reference",
        "desc": "SILVA-derived rRNA refs shipped with SortMeRNA. Extracts ~1.5% of "
                "raw WGS reads as 16S rRNA reads.",
        "tags": [("P3 extract", C_P3L, C_P3BG, C_P3LINE)],
    },
]

DB_ITEM_W = (DB_CARD_W - 0.3) / 5

for i, db in enumerate(DATABASES):
    il = DB_CARD_L + 0.15 + i * DB_ITEM_W
    it = DB_CARD_T + 0.1

    # Vertical divider (between items)
    if i > 0:
        add_rect(slide, il - 0.06, it + 0.05, 0.007, DB_CARD_H - 0.25, fill=C_BORDER)

    # Icon
    add_text_box(slide, db["icon"], il, it, 0.35, 0.3, font_size=Pt(14))

    # Name
    add_text_box(slide, db["name"], il + 0.38, it + 0.02, DB_ITEM_W - 0.55, 0.25,
                 font_size=Pt(8.5), bold=True, color=C_TEXT)
    # Version
    add_text_box(slide, db["ver"], il + 0.38, it + 0.26, DB_ITEM_W - 0.55, 0.18,
                 font_size=Pt(6.5), color=C_MUTED)
    # Desc
    add_text_box(slide, db["desc"], il, it + 0.44, DB_ITEM_W - 0.2, 0.38,
                 font_size=Pt(7), color=C_MUTED, wrap=True)

    # Tags
    tag_l = il
    tag_t = it + 0.86
    for (tag_text, tag_fg, tag_bg, tag_border) in db["tags"]:
        tag_w = len(tag_text) * 0.065 + 0.1
        add_rounded_rect(slide, tag_l, tag_t, tag_w, 0.19,
                         fill=tag_bg, line=tag_border, line_w=Pt(0.4), radius=0.2)
        add_text_box(slide, tag_text, tag_l, tag_t + 0.01, tag_w, 0.18,
                     font_size=Pt(6), bold=True, color=tag_fg,
                     align=PP_ALIGN.CENTER)
        tag_l += tag_w + 0.06


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 3 — COMPLETED DATASETS
# ═══════════════════════════════════════════════════════════════════════════════
DS_T = 4.78
section_label("Completed Datasets", 0.5, DS_T - 0.17)

DS_CARD_W = 6.09
DS_CARD_H = 2.08

DATASETS = [
    {
        "icon":     "💩",
        "title":    "Fecal QC Dataset",
        "subtitle": "Human stool — internal QC cohort",
        "stat_color": C_GOLD,
        "stat_bg":  C_GOLDBG,
        "stat_ln":  C_GOLDLINE,
        "stats": [("18", "Total Samples"), ("3", "Pipelines Run"), ("2", "Data Types")],
        "desc": "Human fecal microbiome samples processed through all three pipelines, "
                "generating paired 16S amplicon and shotgun metagenomics (MGS) "
                "genus-level profiles for cross-method concordance analysis.",
        "pipe_results": [
            ("P1 — GG2 / Woltka",        C_P1L, C_P1BG, C_P1LINE),
            ("P2 — RefSeq / Kraken2",     C_P2L, C_P2BG, C_P2LINE),
            ("P3 — SortMeRNA / Kraken2",  C_P3L, C_P3BG, C_P3LINE),
        ],
        "col": 0,
    },
    {
        "icon":     "⚗️",
        "title":    "Zymo Mock Community",
        "subtitle": "ZymoBIOMICS D6305 — benchmark standard",
        "stat_color": C_PINK,
        "stat_bg":  C_PINKBG,
        "stat_ln":  C_PINKLINE,
        "stats": [("10", "Total Samples"), ("10", "Known Organisms"), ("9/10", "Best TP (P2/P3)")],
        "desc": "Mock community: 8 bacteria + 2 fungi at known abundances. "
                "6 amplicon + 4 MGS libraries benchmarked for TP/FP/FN "
                "against ZymoBIOMICS ground truth at 0.01% RA cutoff.",
        "pipe_results": [
            ("P1  8/10 TP · 29 FP",   C_P1L, C_P1BG, C_P1LINE),
            ("P2  9/10 TP ·  7 FP ✓", C_P2L, C_P2BG, C_P2LINE),
            ("P3  9/10 TP · 40 FP",   C_P3L, C_P3BG, C_P3LINE),
        ],
        "col": 1,
    }
]

for ds in DATASETS:
    l = 0.5 + ds["col"] * (DS_CARD_W + 0.14)
    t = DS_T

    # Card
    add_rounded_rect(slide, l, t, DS_CARD_W, DS_CARD_H,
                     fill=C_SURFACE, line=C_BORDER, line_w=Pt(0.75))

    # Icon
    add_text_box(slide, ds["icon"], l + 0.18, t + 0.1, 0.4, 0.36, font_size=Pt(18))

    # Title
    add_text_box(slide, ds["title"], l + 0.65, t + 0.1, 3.2, 0.28,
                 font_size=Pt(11.5), bold=True, color=C_TEXT)
    # Subtitle
    add_text_box(slide, ds["subtitle"], l + 0.65, t + 0.36, 3.6, 0.18,
                 font_size=Pt(7.5), color=C_MUTED)

    # Status badge
    badge_l = l + DS_CARD_W - 1.42
    add_rounded_rect(slide, badge_l, t + 0.15, 1.2, 0.22,
                     fill=C_GREEN_BG, line=C_GREEN_OK, line_w=Pt(0.5), radius=0.2)
    add_text_box(slide, "✓  Completed", badge_l, t + 0.165, 1.2, 0.2,
                 font_size=Pt(6.5), bold=True, color=C_GREEN_OK,
                 align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, l + 0.18, t + 0.6, DS_CARD_W - 0.36, 0.008, fill=C_BORDER)

    # Stats boxes
    STAT_W = (DS_CARD_W - 0.36 - 0.16) / 3
    for si, (val, key) in enumerate(ds["stats"]):
        sl = l + 0.18 + si * (STAT_W + 0.08)
        st = t + 0.68
        add_rounded_rect(slide, sl, st, STAT_W, 0.44,
                         fill=ds["stat_bg"], line=ds["stat_ln"], line_w=Pt(0.5))
        add_text_box(slide, val, sl, st + 0.02, STAT_W, 0.28,
                     font_size=Pt(16), bold=True,
                     color=ds["stat_color"], align=PP_ALIGN.CENTER)
        add_text_box(slide, key, sl, st + 0.28, STAT_W, 0.17,
                     font_size=Pt(6), color=C_MUTED, align=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, ds["desc"], l + 0.18, t + 1.17, DS_CARD_W - 0.36, 0.42,
                 font_size=Pt(7.5), color=C_MUTED, wrap=True)

    # Pipeline result tags
    tag_l = l + 0.18
    tag_t = t + 1.62
    for (tag_text, tag_fg, tag_bg, tag_border) in ds["pipe_results"]:
        tag_w = len(tag_text) * 0.066 + 0.1
        add_rounded_rect(slide, tag_l, tag_t, tag_w, 0.21,
                         fill=tag_bg, line=tag_border, line_w=Pt(0.4), radius=0.2)
        add_text_box(slide, tag_text, tag_l, tag_t + 0.01, tag_w, 0.2,
                     font_size=Pt(6.5), bold=True, color=tag_fg,
                     align=PP_ALIGN.CENTER)
        tag_l += tag_w + 0.08


# ─── FOOTER BAR ──────────────────────────────────────────────────────────────
F_T = 7.18
add_rounded_rect(slide, 0.5, F_T, 12.33, 0.22,
                 fill=C_SURFACE2, line=C_BORDER, line_w=Pt(0.5))

footer_items = [
    (C_P1, "Pipeline 1: Greengenes2 / Woltka / GTDB taxonomy",   0.7),
    (C_P2, "Pipeline 2: RefSeq full / Kraken2 + Bracken / NCBI", 4.7),
    (C_P3, "Pipeline 3: SortMeRNA 16S extract → Kraken2",         8.5),
]
for color, text, fl in footer_items:
    add_text_box(slide, "●", fl, F_T + 0.035, 0.15, 0.18,
                 font_size=Pt(7), color=color)
    add_text_box(slide, text, fl + 0.15, F_T + 0.035, 3.6, 0.18,
                 font_size=Pt(7), color=C_MUTED)

add_text_box(slide, "github.com/gio9024/16S-and-SM-data-integration",
             10.4, F_T + 0.035, 2.2, 0.18,
             font_size=Pt(6.5), color=C_P1, align=PP_ALIGN.RIGHT)


# ─── SAVE ────────────────────────────────────────────────────────────────────
OUT = "summary_slide.pptx"
prs.save(OUT)
print(f"✅  Saved: {OUT}")
